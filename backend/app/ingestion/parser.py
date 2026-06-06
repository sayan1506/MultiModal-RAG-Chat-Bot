import fitz
from pptx import Presentation
from pathlib import Path
from dataclasses import dataclass


@dataclass
class ParsedPage:
    page_number: int
    text: str
    source_file: str


def parse_pdf(file_path: str) -> list[ParsedPage]:
    doc = fitz.open(file_path)
    pages = []

    for i, page in enumerate(doc):
        text = page.get_text("text")

        pages.append(
            ParsedPage(
                page_number=i + 1,
                text=text.strip(),
                source_file=Path(file_path).name,
            )
        )

    return pages


def parse_pptx(file_path: str) -> list[ParsedPage]:
    prs = Presentation(file_path)
    pages = []

    for i, slide in enumerate(prs.slides):
        text_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_parts.append(shape.text)

        pages.append(
            ParsedPage(
                page_number=i + 1,
                text="\n".join(text_parts).strip(),
                source_file=Path(file_path).name,
            )
        )

    return pages


def parse_file(file_path: str):
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        return parse_pdf(file_path)

    elif ext in (".pptx", ".ppt"):
        return parse_pptx(file_path)

    raise ValueError(f"Unsupported file type: {ext}")